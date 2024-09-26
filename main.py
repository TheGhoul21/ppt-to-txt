from dotenv import load_dotenv
load_dotenv()

import os
import io
import time
from abc import ABC, abstractmethod
from pptx import Presentation
from PIL import Image, ImageDraw, ImageFont
import google.generativeai as genai

class AIModelInterface(ABC):
    @abstractmethod
    def upload_file(self, file_path, mime_type):
        pass

    @abstractmethod
    def wait_for_files_active(self, files):
        pass

    @abstractmethod
    def process_content(self, files, text_content):
        pass

class GeminiModel(AIModelInterface):
    def __init__(self, api_key):
        genai.configure(api_key=api_key)
        self.model = genai.GenerativeModel(
            model_name="gemini-1.5-pro-002",
            generation_config={
                "temperature": 1,
                "top_p": 0.95,
                "top_k": 40,
                "max_output_tokens": 8192,
                "response_mime_type": "text/plain",
            }
        )

    def upload_file(self, file_path, mime_type):
        file = genai.upload_file(file_path, mime_type=mime_type)
        print(f"Uploaded file '{file.display_name}' as: {file.uri}")
        return file

    def wait_for_files_active(self, files):
        print("Waiting for file processing...")
        for name in (file.name for file in files):
            file = genai.get_file(name)
            while file.state.name == "PROCESSING":
                print(".", end="", flush=True)
                time.sleep(10)
                file = genai.get_file(name)
            if file.state.name != "ACTIVE":
                raise Exception(f"File {file.name} failed to process")
        print("...all files ready")
        print()

    def process_content(self, files, text_content):
        chat_history = [{"role": "user", "parts": [file]} for file in files]
        chat_history.append({
            "role": "user",
            "parts": [text_content]
        })
        chat_session = self.model.start_chat(history=chat_history)
        response = chat_session.send_message("Analyze the provided content as instructed.")
        return response.text

def analyze_powerpoint(pptx_file, output_file, ai_model, combine_images=True, add_labels=True):
    prs = Presentation(pptx_file)
    results = []

    def extract_slide_content(slide):
        text = ""
        images = []
        for shape in slide.shapes:
            if hasattr(shape, 'text'):
                text += shape.text + "\n"
            if shape.shape_type == 13:  # Picture
                image = shape.image
                image_bytes = image.blob
                img = Image.open(io.BytesIO(image_bytes))
                images.append(img)
        return text.strip(), images

    def add_id_to_image(img, id_text):
        if not add_labels:
            return img
        draw = ImageDraw.Draw(img)
        font = ImageFont.load_default()
        
        left, top, right, bottom = draw.textbbox((0, 0), id_text, font=font)
        text_width = right - left
        text_height = bottom - top
        
        position = (10, 10)  # Top-left corner, with a small margin
        
        text_bg = Image.new('RGBA', (text_width + 20, text_height + 20), (0, 0, 0, 128))
        img.paste(text_bg, position, text_bg)
        
        draw.text((position[0] + 10, position[1] + 10), id_text, font=font, fill=(255, 255, 255, 255))
        return img
    for i in range(len(prs.slides) - 1):
        slide1 = prs.slides[i]

        text1, images1 = extract_slide_content(slide1)

        combined_text = f"Slide {i + 1}:\n{text1}"
        combined_images = images1

        processed_imgs = []
        ai_files = []

        if combined_images:
            for idx, img in enumerate(combined_images):
                id_text = f"IMG_{i+1}_{idx+1}"
                img_with_id = add_id_to_image(img.copy(), id_text)
                processed_imgs.append(img_with_id)
                
                img_path = f"temp_image_{i+1}_{idx+1}.png"
                img_with_id.save(img_path)
                
                ai_file = ai_model.upload_file(img_path, mime_type="image/png")
                ai_files.append(ai_file)
                os.remove(img_path)  # Remove temporary file

            if combine_images:
                widths, heights = zip(*(i.size for i in processed_imgs))
                max_height = max(heights)
                total_width = sum(widths)
                combined_img = Image.new('RGB', (total_width, max_height))
                x_offset = 0
                for img in processed_imgs:
                    combined_img.paste(img, (x_offset, 0))
                    x_offset += img.size[0]
                
                combined_img_path = f"combined_image_{i}.png"
                combined_img.save(combined_img_path)
                
                ai_files = [ai_model.upload_file(combined_img_path, mime_type="image/png")]
                os.remove(combined_img_path)  # Remove temporary file

        # Wait for files to be processed
        ai_model.wait_for_files_active(ai_files)

        # Prepare text content
        analysis_instructions = (
            "Analyze the provided slides and images. Create a cohesive summary of the content, "
            "integrating relevant information from both the text and images. For each image that "
            "adds value to the summary:\n"
            "1. Reference the image using its ID (e.g., IMG_1_1)\n"
            "2. Provide a brief, descriptive caption for the image\n"
            "3. Place the image reference immediately after the relevant text\n\n"
            "Format image references as: [image src=IMG_X_X.png caption=\"Your caption here\"]\n\n"
            "Ignore any images that don't contribute significant information to the summary. "
            "Focus on creating a flowing, informative summary that incorporates text and "
            "image references seamlessly.\n\n"
            f"Here's the text content from the slides:\n{combined_text}"
        )

        # Process content with AI model
        result = ai_model.process_content(ai_files, analysis_instructions)
        results.append(result)

    # Write results to output file
    with open(output_file, 'w', encoding='utf-8') as f:
        for result in results:
            f.write(result + "\n\n")

# Initialize the Gemini model
gemini_model = GeminiModel(api_key=os.environ["GEMINI_API_KEY"])

pptx_file = "files/fivhterDeckSA.pptx"
output_file = "fivhterDeckSA.txt"
analyze_powerpoint(pptx_file, output_file, gemini_model, combine_images=True, add_labels=True)
